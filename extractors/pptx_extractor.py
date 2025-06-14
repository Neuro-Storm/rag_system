from pptx import Presentation

def extract_text_from_pptx(file_path: str) -> str:
    """Извлечение текста из PPTX"""
    try:
        prs = Presentation(file_path)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return "\n".join(text)
    except Exception as e:
        print(f"Ошибка чтения PPTX: {str(e)}")
        return ""